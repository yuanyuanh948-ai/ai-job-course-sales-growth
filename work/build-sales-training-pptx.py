from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt

ROOT = Path.cwd()
OUT = ROOT / "outputs" / "AI就业班产品培训PPT-销售团队版.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

INK = RGBColor(29, 29, 31)
MUTED = RGBColor(110, 110, 115)
BLUE = RGBColor(0, 113, 227)
LIGHT = RGBColor(245, 245, 247)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)


def set_text(box, text, size=24, color=INK, bold=False, align=PP_ALIGN.LEFT):
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei UI"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def add_title(slide, kicker, title, subtitle=None):
    k = slide.shapes.add_textbox(Inches(0.65), Inches(0.42), Inches(5.2), Inches(0.35))
    set_text(k, kicker.upper(), 9, BLUE, True)
    t = slide.shapes.add_textbox(Inches(0.62), Inches(0.82), Inches(9.8), Inches(1.0))
    set_text(t, title, 30, INK, False)
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.66), Inches(1.75), Inches(9.6), Inches(0.45))
        set_text(s, subtitle, 14, MUTED)


def add_footer(slide, n, total=18):
    f = slide.shapes.add_textbox(Inches(11.65), Inches(6.95), Inches(1.0), Inches(0.24))
    set_text(f, f"{n:02d} / {total}", 8, MUTED, False, PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, fill=LIGHT, title_color=INK, body_color=MUTED):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tx = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.18), Inches(w - 0.36), Inches(0.36))
    set_text(tx, title, 12, title_color, True)
    bx = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.62), Inches(w - 0.36), Inches(h - 0.75))
    set_text(bx, body, 13, body_color)


def blank(bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


slides = [
    ("cover", "AI就业班产品培训", "让销售听懂：模块作用、学生收获、服务交付", []),
    ("培训目标", "今天不是学技术，而是学会讲清产品价值", "", ["听懂产品路径", "会讲模块价值", "避免结果型承诺"]),
    ("一句话价值", "看懂岗位 → 补齐能力 → 做出项目 → 整理材料 → 进入求职准备", "", ["产品不是课程表，而是一条求职准备路径"]),
    ("用户痛点", "销售先讲痛点", "", ["方向不清", "基础不稳", "项目太弱", "表达不深"]),
    ("目标人群", "三类适配人群", "", ["计算机学生：项目弱、简历空", "互联网转岗：经验迁移到 AI", "跨专业人群：需要清晰路径", "不适配：不投入、不写代码、只期待结果承诺"]),
    ("产品路径", "六个模块，一条主线", "", ["新手通识", "Python基础与实战", "机器学习与深度学习", "Vibe Coding", "就业技术栈", "综合项目与作品集"]),
    ("模块1", "新手通识", "让学生从迷茫进入可执行状态", ["解决方向迷茫", "解决环境卡点", "跑通项目入口"]),
    ("模块2", "Python基础与实战", "补齐真实项目的底层执行能力", ["文件与数据处理", "清洗与可视化", "电商分析报告生成器"]),
    ("模块3", "机器学习与深度学习", "从工具层表达升级到模型层理解", ["ML流水线", "模型调优", "PyTorch与深度学习项目"]),
    ("模块4", "Vibe Coding", "借助 AI 提升项目开发效率", ["读代码", "改代码", "定位报错", "拆任务推进"]),
    ("模块5", "就业技术栈", "把 AI 能力转成就业作品", ["RAG + Prompt", "Agent 工作流", "Skill 封装", "MCP / Hooks / Harness"]),
    ("模块6", "综合项目与作品集", "把学习结果变成求职材料", ["Demo", "README", "复盘", "简历项目描述"]),
    ("服务价值", "服务不是赠品，是交付系统", "", ["入学启动", "过程监督", "答疑支持", "项目推进", "求职准备", "资料供给"]),
    ("销售方法", "先痛点，再路径，最后成果", "", ["先讲痛点", "再讲路径", "讲可见产出", "讲服务保障", "讲边界"]),
    ("不要这样讲", "三个禁区", "", ["不堆技术名词", "不只说课很多", "不做结果型承诺"]),
    ("用户讲法", "三类用户，三种重点", "", ["学生：项目和面试", "转岗：经验迁移", "跨专业：路径和门槛"]),
    ("成交模板", "四句话讲完整", "", ["一句话总价值", "模块价值", "服务保障", "可见产出"]),
    ("结尾", "路径 · 项目 · 表达", "销售不需要讲懂每节课，但要讲清学生最后能得到什么。", []),
]

for idx, (kicker, title, subtitle, items) in enumerate(slides, start=1):
    if idx == 1:
        s = blank(BLUE)
        tb = s.shapes.add_textbox(Inches(0.72), Inches(2.2), Inches(9.8), Inches(1.4))
        set_text(tb, title, 44, WHITE)
        st = s.shapes.add_textbox(Inches(0.78), Inches(3.6), Inches(8.6), Inches(0.5))
        set_text(st, subtitle, 18, WHITE)
        add_footer(s, idx)
        continue
    s = blank(WHITE if idx % 3 else LIGHT)
    add_title(s, kicker, title, subtitle)
    if items:
        cols = 3 if len(items) in (3, 6) else 2
        w = 3.85 if cols == 3 else 5.75
        h = 1.35
        start_x = 0.75
        start_y = 2.55
        for i, item in enumerate(items):
            x = start_x + (i % cols) * (w + 0.35)
            y = start_y + (i // cols) * (h + 0.35)
            fill = BLUE if ("就业技术栈" in item or i == len(items) - 1 and idx in (3, 14, 18)) else LIGHT
            title_color = WHITE if fill == BLUE else INK
            body_color = WHITE if fill == BLUE else MUTED
            if "：" in item:
                ct, cb = item.split("：", 1)
            else:
                ct, cb = item, ""
            add_card(s, x, y, w, h, ct, cb or item, fill, title_color, body_color)
    add_footer(s, idx)

prs.save(OUT)
print(OUT)
